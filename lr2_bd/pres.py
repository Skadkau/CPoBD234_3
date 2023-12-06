import pptx

def replace_keywords_in_docx(docx_path, keyword_dict, output_path):
    doc = pptx.Presentation(docx_path)
    slovar = keyword_dict
    paragraph = doc.slides[0]
    for elem in paragraph.placeholders:
        for k, v in slovar.items():
            if k in elem.text:
                elem.text = elem.text.replace(k, v)

    doc.save(output_path)

# Пример использования
keyword_dict = {
    "{user}": "Иван Иванов",
    "{skill}": "Значение2",
    "{add skill}": "Значение3",
    "{datta}": "Значение4",
}
#файл с ключами
input_file = "Sertifikat v2.pptx"
# название итогового файта
output_file = "Sertifikat v2.pptx"

replace_keywords_in_docx(input_file, keyword_dict, output_file)